import collections
import collections.abc
collections.Container = collections.abc.Container
collections.Mapping = collections.abc.Mapping
collections.MutableMapping = collections.abc.MutableMapping
collections.Sequence = collections.abc.Sequence
collections.Set = collections.abc.Set

import sys
from pptx import Presentation
from docx import Document
import pypandoc
import re
import os

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text_runs = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_runs.append("\n".join(slide_text))
    return text_runs

def clean_text(text):
    # Removing vertical tab character \x0b
    text = text.replace('\x0b', '')

    # Detecting non-allowed characters
    allowed_chars_regex = r'[^\t\n\r\u0020-\uD7FF\uE000-\uFFFD\u10000-\u10FFFF]'
    problematic_chars = re.findall(allowed_chars_regex, text)
    if problematic_chars:
        print(f"Problematic characters found: {problematic_chars}")

    # Removing problematic characters
    cleaned_text = re.sub(allowed_chars_regex, '', text)
    return cleaned_text

def save_text_to_docx(text_runs, docx_path):
    doc = Document()
    for idx, slide_text in enumerate(text_runs, 1):
        doc.add_heading(f'Slide {idx}', level=1)
        doc.add_paragraph(clean_text(slide_text))
    doc.save(docx_path)

def convert_docx_to_rtf(docx_path, rtf_path):
    pypandoc.convert_file(docx_path, 'rtf', outputfile=rtf_path)

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python pptx_2_docx.py <input_pptx_file>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    base_name = os.path.splitext(pptx_path)[0]
    docx_path = f"{base_name}.docx"
    rtf_path = f"{base_name}.rtf"

    text_runs = extract_text_from_pptx(pptx_path)
    save_text_to_docx(text_runs, docx_path)
    convert_docx_to_rtf(docx_path, rtf_path)

    print(f'Text extracted and saved to {docx_path} and {rtf_path}')
